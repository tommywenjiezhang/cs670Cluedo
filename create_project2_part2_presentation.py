#!/usr/bin/env python3
# Install with: pip install python-pptx

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import List

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx is not installed.")
    print("Install it with: pip install python-pptx")
    sys.exit(1)


INPUT_PATH = Path("project2_part2_formal_presentation.md")
OUTPUT_PATH = Path("Project2_Part2_Cluedo_Formal_Presentation.pptx")

PROJECT_FOOTER = "Project 2 Part 2 - Cluedo Digital Edition"
FONT_NAME = "Calibri"
TITLE_FONT_SIZE = Pt(38)
SLIDE_TITLE_FONT_SIZE = Pt(32)
BULLET_FONT_SIZE = Pt(22)
SUB_BULLET_FONT_SIZE = Pt(18)
SUBTITLE_FONT_SIZE = Pt(22)
BODY_MAX_LINES = 12
BODY_MAX_BULLETS = 6
DEMO_MAX_STEPS = 3


@dataclass
class ContentItem:
    kind: str
    text: str
    level: int = 0


@dataclass
class SlideData:
    raw_title: str
    title: str
    items: List[ContentItem] = field(default_factory=list)
    notes: List[str] = field(default_factory=list)
    is_title_slide: bool = False


def clean_markdown(text: str) -> str:
    text = text.replace("**", "")
    text = text.replace("`", "")
    text = re.sub(r"\[(.*?)\]\((.*?)\)", r"\1", text)
    return re.sub(r"\s+", " ", text).strip()


def extract_heading_title(line: str) -> str | None:
    match = re.match(r"^##\s+Slide\s+\d+\s*:\s*(.+?)\s*$", line)
    if not match:
        return None
    return clean_markdown(match.group(1))


def parse_markdown(markdown_text: str) -> List[SlideData]:
    slides: List[SlideData] = []
    current: SlideData | None = None
    mode = "body"

    for raw_line in markdown_text.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()

        title = extract_heading_title(stripped)
        if title is not None:
            current = SlideData(
                raw_title=title,
                title=title,
                is_title_slide=(title.lower() == "title"),
            )
            slides.append(current)
            mode = "body"
            continue

        if current is None or not stripped:
            continue

        if stripped.startswith("**Speaker Notes:**"):
            mode = "notes"
            note_body = clean_markdown(stripped[len("**Speaker Notes:**"):])
            if note_body:
                current.notes.append(note_body)
            continue

        if stripped.startswith("**Demo Cue:**"):
            mode = "demo_cue"
            cue_body = clean_markdown(stripped[len("**Demo Cue:**"):])
            if cue_body:
                current.notes.append(f"Demo Cue: {cue_body}")
            continue

        if mode in {"notes", "demo_cue"} and not stripped.startswith(("-", "*")) and not re.match(r"^\d+\.\s+", stripped):
            text = clean_markdown(stripped)
            if text:
                if mode == "demo_cue":
                    current.notes.append(f"Demo Cue: {text}")
                else:
                    current.notes.append(text)
            continue

        mode = "body"

        bullet_match = re.match(r"^(\s*)[-*]\s+(.+)$", line)
        number_match = re.match(r"^(\s*)(\d+)\.\s+(.+)$", line)
        colon_line_match = re.match(r"^(\s*)([A-Za-z][A-Za-z ]+:\s+.+)$", line)

        if bullet_match:
            indent = len(bullet_match.group(1).replace("\t", "    "))
            level = 1 if indent >= 2 else 0
            current.items.append(ContentItem("bullet", clean_markdown(bullet_match.group(2)), level))
            continue

        if number_match:
            indent = len(number_match.group(1).replace("\t", "    "))
            level = 1 if indent >= 2 else 0
            current.items.append(ContentItem("number", clean_markdown(number_match.group(3)), level))
            continue

        if colon_line_match:
            indent = len(colon_line_match.group(1).replace("\t", "    "))
            level = 1 if indent >= 2 else 0
            current.items.append(ContentItem("bullet", clean_markdown(colon_line_match.group(2)), level))
            continue

    return slides


def estimate_lines(items: List[ContentItem]) -> int:
    total = 0
    for item in items:
        total += 1
        total += max(0, len(item.text) // 70)
    return total


def build_continuation_title(base_title: str, index: int) -> str:
    return f"{base_title} Continued" if index == 1 else f"{base_title} Continued {index}"


def split_regular_items(slide: SlideData) -> List[List[ContentItem]]:
    chunks: List[List[ContentItem]] = []
    current: List[ContentItem] = []
    top_level_count = 0
    line_count = 0

    for item in slide.items:
        item_lines = 1 + max(0, len(item.text) // 70)
        prospective_top_level = top_level_count + (1 if item.level == 0 else 0)
        if current and (
            prospective_top_level > BODY_MAX_BULLETS
            or line_count + item_lines > BODY_MAX_LINES
        ):
            chunks.append(current)
            current = []
            top_level_count = 0
            line_count = 0

        current.append(item)
        if item.level == 0:
            top_level_count += 1
        line_count += item_lines

    if current:
        chunks.append(current)

    return chunks


def split_demo_items(slide: SlideData) -> List[List[ContentItem]]:
    groups: List[List[ContentItem]] = []
    current_group: List[ContentItem] = []
    for item in slide.items:
        if item.kind == "number" and item.level == 0:
            if current_group:
                groups.append(current_group)
            current_group = [item]
        else:
            if not current_group:
                current_group = [item]
            else:
                current_group.append(item)
    if current_group:
        groups.append(current_group)

    chunks: List[List[ContentItem]] = []
    for start in range(0, len(groups), DEMO_MAX_STEPS):
        chunk: List[ContentItem] = []
        for group in groups[start:start + DEMO_MAX_STEPS]:
            chunk.extend(group)
        chunks.append(chunk)
    return chunks


def split_slide(slide: SlideData) -> List[SlideData]:
    if slide.is_title_slide:
        return [slide]

    has_numbered_steps = any(item.kind == "number" and item.level == 0 for item in slide.items)
    if has_numbered_steps:
        item_chunks = split_demo_items(slide)
    else:
        item_chunks = split_regular_items(slide)

    if len(item_chunks) <= 1:
        return [slide]

    result: List[SlideData] = []
    for index, chunk in enumerate(item_chunks):
        title = slide.title if index == 0 else build_continuation_title(slide.title, index)
        result.append(
            SlideData(
                raw_title=slide.raw_title,
                title=title,
                items=chunk,
                notes=list(slide.notes),
                is_title_slide=False,
            )
        )
    return result


def set_run_style(run, font_size, bold=False):
    run.font.name = FONT_NAME
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(40, 40, 40)


def configure_text_frame(text_frame):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = Inches(0.08)
    text_frame.margin_right = Inches(0.08)
    text_frame.margin_top = Inches(0.04)
    text_frame.margin_bottom = Inches(0.04)


def add_footer(slide, slide_number: int):
    footer_box = slide.shapes.add_textbox(Inches(0.45), Inches(6.85), Inches(6.2), Inches(0.28))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    footer_para = footer_frame.paragraphs[0]
    footer_run = footer_para.add_run()
    footer_run.text = PROJECT_FOOTER
    set_run_style(footer_run, Pt(10))
    footer_run.font.color.rgb = RGBColor(100, 100, 100)

    number_box = slide.shapes.add_textbox(Inches(11.9), Inches(6.82), Inches(0.45), Inches(0.28))
    number_frame = number_box.text_frame
    number_frame.clear()
    number_para = number_frame.paragraphs[0]
    number_para.alignment = PP_ALIGN.RIGHT
    number_run = number_para.add_run()
    number_run.text = str(slide_number)
    set_run_style(number_run, Pt(10))
    number_run.font.color.rgb = RGBColor(100, 100, 100)


def set_slide_notes(slide, notes: List[str]):
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None or not hasattr(notes_slide, "notes_text_frame"):
        raise RuntimeError(
            "The installed python-pptx version does not support writing speaker notes."
        )

    notes_text = "\n\n".join(note for note in notes if note.strip())
    notes_frame = notes_slide.notes_text_frame
    notes_frame.clear()
    paragraph = notes_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = notes_text


def create_title_slide(prs: Presentation, slide_data: SlideData):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    visible_lines = [item.text for item in slide_data.items]
    title_text = visible_lines[0] if visible_lines else "Project 2 Part 2"
    subtitle_lines = visible_lines[1:]

    title_shape.text = title_text
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.LEFT
    for run in title_para.runs:
        set_run_style(run, TITLE_FONT_SIZE, bold=True)

    subtitle_shape.text = "\n".join(subtitle_lines)
    for para in subtitle_shape.text_frame.paragraphs:
        for run in para.runs:
            set_run_style(run, SUBTITLE_FONT_SIZE)

    if slide_data.notes:
        set_slide_notes(slide, slide_data.notes)


def find_body_placeholder(slide):
    for shape in slide.placeholders:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return shape
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape != slide.shapes.title:
            return shape
    raise RuntimeError("No body placeholder found for slide layout.")


def create_content_slide(prs: Presentation, slide_data: SlideData, slide_number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = slide_data.title
    title_para = title_shape.text_frame.paragraphs[0]
    for run in title_para.runs:
        set_run_style(run, SLIDE_TITLE_FONT_SIZE, bold=True)

    body_shape = find_body_placeholder(slide)
    text_frame = body_shape.text_frame
    configure_text_frame(text_frame)

    first = True
    for item in slide_data.items:
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        paragraph.level = item.level
        paragraph.space_after = Pt(6 if item.level == 0 else 2)
        paragraph.space_before = Pt(0)

        prefix = ""
        if item.kind == "number" and item.level == 0:
            prefix = "• "
        run = paragraph.add_run()
        run.text = f"{prefix}{item.text}"
        set_run_style(run, BULLET_FONT_SIZE if item.level == 0 else SUB_BULLET_FONT_SIZE)

    add_footer(slide, slide_number)

    if slide_data.notes:
        set_slide_notes(slide, slide_data.notes)


def create_presentation(slides: List[SlideData]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    visible_slide_number = 0
    for slide_data in slides:
        if slide_data.is_title_slide:
            create_title_slide(prs, slide_data)
        else:
            visible_slide_number += 1
            create_content_slide(prs, slide_data, visible_slide_number)
    return prs


def main() -> int:
    print(f"Reading {INPUT_PATH.name}...")
    if not INPUT_PATH.exists():
        print(f"Error: {INPUT_PATH.name} does not exist.")
        return 1

    markdown_text = INPUT_PATH.read_text(encoding="utf-8")
    if not markdown_text.strip():
        print(f"Error: {INPUT_PATH.name} is empty.")
        return 1

    print("Parsing slides...")
    parsed_slides = parse_markdown(markdown_text)
    if not parsed_slides:
        print("Error: No slides were found. Expected headings like '## Slide X: ...'.")
        return 1

    expanded_slides: List[SlideData] = []
    for slide in parsed_slides:
        expanded_slides.extend(split_slide(slide))

    print("Creating PowerPoint presentation...")
    try:
        presentation = create_presentation(expanded_slides)
    except Exception as exc:
        print(f"Error while creating presentation: {exc}")
        return 1

    print(f"Saving {OUTPUT_PATH.name}...")
    try:
        presentation.save(str(OUTPUT_PATH))
    except Exception as exc:
        print(f"Error while saving PowerPoint file: {exc}")
        return 1

    if not OUTPUT_PATH.exists() or OUTPUT_PATH.stat().st_size == 0:
        print("Error: PowerPoint file was not created successfully.")
        return 1

    print("Done.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
